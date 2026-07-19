from __future__ import annotations

from pathlib import Path
from zipfile import ZIP_DEFLATED, ZIP_STORED, ZipFile

import pytest

from core.worker_entry import convert_one


def make_pdf(path: Path, marker: str) -> None:
    stream = f"BT /F1 18 Tf 72 720 Td ({marker}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n"
        + stream
        + b"\nendstream",
    ]
    document = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, body in enumerate(objects, start=1):
        offsets.append(len(document))
        document.extend(f"{number} 0 obj\n".encode("ascii"))
        document.extend(body)
        document.extend(b"\nendobj\n")
    xref_offset = len(document)
    document.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    document.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        document.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    document.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_offset}\n%%EOF\n".encode("ascii")
    )
    path.write_bytes(document)


def make_docx(path: Path, marker: str) -> None:
    content_types = """<?xml version="1.0" encoding="UTF-8"?>
    <Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
      <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
      <Default Extension="xml" ContentType="application/xml"/>
      <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
    </Types>"""
    relationships = """<?xml version="1.0" encoding="UTF-8"?>
    <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
      <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
    </Relationships>"""
    document = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
    <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
      <w:body><w:p><w:r><w:t>{marker}</w:t></w:r></w:p><w:sectPr/></w:body>
    </w:document>"""
    with ZipFile(path, "w", ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", relationships)
        archive.writestr("word/document.xml", document)


def make_epub(path: Path, marker: str) -> None:
    container = """<?xml version="1.0"?>
    <container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
      <rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/></rootfiles>
    </container>"""
    package = """<?xml version="1.0" encoding="UTF-8"?>
    <package xmlns="http://www.idpf.org/2007/opf" version="2.0" unique-identifier="id">
      <metadata xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>Test</dc:title><dc:identifier id="id">test</dc:identifier></metadata>
      <manifest><item id="chapter" href="chapter.xhtml" media-type="application/xhtml+xml"/></manifest>
      <spine><itemref idref="chapter"/></spine>
    </package>"""
    chapter = f"""<html xmlns="http://www.w3.org/1999/xhtml"><head><title>Test</title></head>
    <body><h1>{marker}</h1><p>EPUB paragraph</p></body></html>"""
    with ZipFile(path, "w") as archive:
        archive.writestr("mimetype", "application/epub+zip", compress_type=ZIP_STORED)
        archive.writestr("META-INF/container.xml", container, compress_type=ZIP_DEFLATED)
        archive.writestr("OEBPS/content.opf", package, compress_type=ZIP_DEFLATED)
        archive.writestr("OEBPS/chapter.xhtml", chapter, compress_type=ZIP_DEFLATED)


def create_format_samples(directory: Path) -> dict[str, tuple[Path, str]]:
    samples: dict[str, tuple[Path, str]] = {}

    simple_text = {
        "txt": "TEXT_MARKER",
        "html": "HTML_MARKER",
        "csv": "CSV_MARKER",
        "json": "JSON_MARKER",
        "xml": "XML_MARKER",
        "rtf": "RTF_MARKER",
    }
    bodies = {
        "txt": "TEXT_MARKER\n",
        "html": "<html><body><h1>HTML_MARKER</h1></body></html>",
        "csv": "name,value\nCSV_MARKER,1\n",
        "json": '{"name": "JSON_MARKER"}',
        "xml": "<root><name>XML_MARKER</name></root>",
        "rtf": r"{\rtf1\ansi RTF_MARKER}",
    }
    for extension, marker in simple_text.items():
        path = directory / f"sample.{extension}"
        path.write_text(bodies[extension], encoding="utf-8")
        samples[extension] = (path, marker)

    pdf = directory / "sample.pdf"
    make_pdf(pdf, "PDF_MARKER")
    samples["pdf"] = (pdf, "PDF_MARKER")

    docx = directory / "sample.docx"
    make_docx(docx, "DOCX_MARKER")
    samples["docx"] = (docx, "DOCX_MARKER")

    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "PPTX_MARKER"
    pptx = directory / "sample.pptx"
    presentation.save(pptx)
    samples["pptx"] = (pptx, "PPTX_MARKER")

    from openpyxl import Workbook

    workbook = Workbook()
    workbook.active["A1"] = "XLSX_MARKER"
    xlsx = directory / "sample.xlsx"
    workbook.save(xlsx)
    samples["xlsx"] = (xlsx, "XLSX_MARKER")

    import xlwt

    legacy_workbook = xlwt.Workbook()
    sheet = legacy_workbook.add_sheet("Sheet1")
    sheet.write(0, 0, "XLS_MARKER")
    xls = directory / "sample.xls"
    legacy_workbook.save(str(xls))
    samples["xls"] = (xls, "XLS_MARKER")

    epub = directory / "sample.epub"
    make_epub(epub, "EPUB_MARKER")
    samples["epub"] = (epub, "EPUB_MARKER")

    zip_path = directory / "sample.zip"
    with ZipFile(zip_path, "w", ZIP_DEFLATED) as archive:
        archive.writestr("inside.txt", "ZIP_MARKER")
    samples["zip"] = (zip_path, "ZIP_MARKER")
    return samples


@pytest.mark.parametrize(
    "extension",
    ["txt", "html", "csv", "json", "xml", "rtf", "pdf", "docx", "pptx", "xlsx", "xls", "epub", "zip"],
)
def test_real_supported_format_conversion(tmp_path: Path, extension: str) -> None:
    source, marker = create_format_samples(tmp_path)[extension]
    destination = tmp_path / f"result-{extension}.md"

    convert_one(source, destination)

    result = destination.read_text(encoding="utf-8").replace("\\_", "_")
    assert marker in result


def test_corrupt_supported_document_fails_without_partial_output(tmp_path: Path) -> None:
    source = tmp_path / "corrupt.pdf"
    destination = tmp_path / "corrupt.md"
    source.write_bytes(b"not a real PDF")

    with pytest.raises(Exception):
        convert_one(source, destination)

    assert not destination.exists()
